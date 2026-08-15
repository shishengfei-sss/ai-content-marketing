#!/usr/bin/env python3
"""移除 PPT 右下角页码（如 01 / 27）。"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
PPT = ROOT / "docs/07-对外资料/03-汇报材料/内容获客平台_全域四平台实施方案_20260804.pptx"

PAGE_NUM = re.compile(r"^\d{1,2}\s*/\s*\d+$")


def remove_page_numbers(prs: Presentation) -> int:
    n = 0
    for slide in prs.slides:
        for sh in list(slide.shapes):
            if hasattr(sh, "text") and PAGE_NUM.fullmatch(sh.text.strip()):
                el = sh.element
                el.getparent().remove(el)
                n += 1
    return n


def main() -> None:
    prs = Presentation(str(PPT))
    count = remove_page_numbers(prs)
    try:
        prs.save(str(PPT))
        print(f"已移除 {count} 处页码：{PPT}")
    except PermissionError:
        alt = PPT.with_name(PPT.stem + "_无页码.pptx")
        prs.save(str(alt))
        print(f"原文件被占用，已另存（{count} 处页码已移除）：{alt}")


if __name__ == "__main__":
    main()
