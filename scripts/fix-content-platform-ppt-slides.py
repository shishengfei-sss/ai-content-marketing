#!/usr/bin/env python3
"""精修 PPT 第 13～14 页：按 shape 坐标填字，清除素材页残留。"""

from __future__ import annotations

import shutil
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
DIR = ROOT / "docs/07-对外资料/03-汇报材料"
SRC = DIR / "内容获客平台_全域四平台实施方案_20260804_backup_integrate_20260805_142920.pptx"
OUT = DIR / "内容获客平台_全域四平台实施方案_20260804.pptx"
FONT = "微软雅黑"
FOOTER = "内容获客平台 · AI 内容工厂 + 交易闭环 + CRM"


def write_shape(shape, text: str, *, clear_only: bool = False) -> None:
    if not shape.has_text_frame:
        return
    ref = None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            ref = run
            break
        if ref:
            break
    size = ref.font.size if ref and ref.font.size else Pt(12)
    bold = ref.font.bold if ref else False
    color = RGBColor(0x11, 0x18, 0x27)
    if ref:
        try:
            color = ref.font.color.rgb
        except AttributeError:
            pass
    if 140_000 < shape.top < 220_000:
        color = RGBColor(0xFF, 0xFF, 0xFF)
        size = Pt(24)
        bold = True
    elif 450_000 < shape.top < 650_000:
        color = RGBColor(0x38, 0xBD, 0xF8)
        size = Pt(12)
        bold = False
    tf = shape.text_frame
    tf.clear()
    content = "" if clear_only else text
    if not content:
        return
    for i, line in enumerate(content.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.name = FONT
            run.font.size = size
            run.font.bold = bold
            try:
                run.font.color.rgb = color
            except AttributeError:
                pass


def shape_at(slide, top: int, left: int, tol: int = 80_000):
    for sh in slide.shapes:
        if abs(sh.top - top) < tol and abs(sh.left - left) < tol and sh.has_text_frame:
            return sh
    return None


def fix_slide13(slide) -> None:
    write_shape(shape_at(slide, 164592, 640080), "交易合规 · 可售商品四道闸")
    write_shape(
        shape_at(slide, 502920, 640080),
        "卖什么要能审 —— 与「素材合规」分工：素材管创作导出，交易管可售 SKU",
    )
    write_shape(shape_at(slide, 1188720, 1005840), "主体资质闸")
    write_shape(
        shape_at(slide, 1188720, 4389120),
        "个人 / 个体工商户 / 企业 分类入驻\n执照或身份证 · 支付进件",
    )
    write_shape(shape_at(slide, 1554480, 4389120), "", clear_only=True)

    write_shape(shape_at(slide, 2240280, 1005840), "商品上架闸")
    write_shape(
        shape_at(slide, 2240280, 4389120),
        "机审 + 平台人审 · 通过后方可私域上架\n未过审不可上架销售",
    )
    write_shape(shape_at(slide, 2606040, 4389120), "", clear_only=True)

    write_shape(shape_at(slide, 3291840, 1005840), "公域与稽查")
    write_shape(
        shape_at(slide, 3291840, 4389120),
        "公域挂载闸：过上架闸后才可同步抖店 / 课程库\n"
        "事后稽查闸：举报抽检 · 强制下架（存量订单照常履约）",
    )
    write_shape(shape_at(slide, 3657600, 4389120), "", clear_only=True)

    write_shape(shape_at(slide, 4160520, 640080), "平台能力（与素材闸并列）")
    write_shape(shape_at(slide, 4526280, 640080), "", clear_only=True)

    write_shape(shape_at(slide, 4572000, 777240), "商家 · 多店")
    write_shape(
        shape_at(slide, 5029200, 822960),
        "入驻 → 多店铺 → 商品挂店下\n配额受当前套餐限制",
    )
    write_shape(shape_at(slide, 4572000, 4480560), "SaaS 套餐")
    write_shape(
        shape_at(slide, 5029200, 4526280),
        "免费 / 基础 / 旗舰 平台可配\n按生效～失效周期管理权益",
    )
    write_shape(shape_at(slide, 4572000, 8183880), "与素材闸")
    write_shape(
        shape_at(slide, 5029200, 8229600),
        "素材：审创作导出物\n交易：审可售卖商品",
    )
    # 清除误放在中间的宽栏
    write_shape(shape_at(slide, 4526280, 4343400), "", clear_only=True)


def fix_slide14(slide) -> None:
    write_shape(shape_at(slide, 502920, 640080), "入驻 → 合规上架 → 付钱 → 履约 → 退款关权益")
    left_body = (
        "· 三类主体入驻（个人/个体/企业）\n"
        "· 多店铺：商品挂店下 · 配额受套餐限制\n"
        "· P10 套餐配置 · P11 权益开通（验收期人工）\n"
        "· 虚拟：专栏/单课、资料包、会员卡\n"
        "· 服务：咨询预约、次数卡、简单核销\n"
        "· 统一支付 + 权益状态机（开通/到期/退款）\n"
        "· 小程序履约端 · 订单/买家 · 抖音回流\n"
        "· 商品合规：机审 + 人审 · 过审才可挂抖店"
    )
    write_shape(shape_at(slide, 1783080, 868680, tol=120_000), left_body)
    right_body = (
        "· 套餐自助购 · 剩余天数折算 → P2\n"
        "· 店铺装修 / 向导 → P2\n"
        "· H5 完整商城 / 混合货架 → P2\n"
        "· 实物快递 → P2；门店核销 → P3\n"
        "· 深度微页面 / PC 网校 → P3\n"
        "· 四平台齐发 → 单公域过后再扩\n"
        "· 永久不做：多级分销、买家=Contact"
    )
    write_shape(shape_at(slide, 1737360, 6446520, tol=120_000), right_body)
    write_shape(
        shape_at(slide, 5669280, 868680, tol=120_000),
        "硬验收：入驻 → 商品合规 → 私域可售 → 公域挂载 → 履约 → 退款关权益。装修开店 Phase 2 补。",
    )


def fix_slide2(slide) -> None:
    sh = shape_at(slide, 0, 0, tol=9_999_999)  # noop fallback
    for s in slide.shapes:
        if s.has_text_frame and "③ 商城" in s.text and "④" in s.text:
            write_shape(
                s,
                "① 平台三层（商家多店）  ② 智能体分期\n"
                "③ 商城首期 + 交易合规  ④ 四平台分阶段挂载",
            )


def fix_slide10(slide) -> None:
    write_shape(
        shape_at(slide, 2240280, 4389120),
        "交易层 · 商家多店 · 套餐 SaaS · 商品/订单/权益",
    )
    write_shape(shape_at(slide, 2606040, 4389120), "Phase 1：入驻过审 · 商品合规 · 能卖能履约")
    sh = None
    for s in slide.shapes:
        if s.has_text_frame and s.text.strip() == "AI 生产内容 + 平台成交":
            sh = s
    if sh:
        write_shape(sh, "AI 生产 + 多店成交 + 套餐可扩展")


def fix_footers(slide) -> None:
    footers = [
        s
        for s in slide.shapes
        if s.has_text_frame and s.top > 6_000_000 and len(s.text.strip()) < 80
    ]
    if not footers:
        return
    primary = min(footers, key=lambda s: (s.top, s.left))
    write_shape(primary, FOOTER)
    for s in footers:
        if s is not primary:
            el = s.element
            el.getparent().remove(el)


def main() -> None:
  # 若 OUT 已存在且页数正确，在其上精修；否则从 SRC 恢复
    if OUT.exists():
        prs_check = Presentation(str(OUT))
        if len(prs_check.slides) == 28:
            prs = prs_check
        else:
            shutil.copy2(SRC, OUT)
            prs = Presentation(str(OUT))
    else:
        shutil.copy2(SRC, OUT)
        prs = Presentation(str(OUT))
    fix_slide13(prs.slides[12])
    fix_slide14(prs.slides[13])
    fix_slide2(prs.slides[1])
    fix_slide10(prs.slides[9])
    for s in prs.slides:
        fix_footers(s)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    backup = OUT.with_name(OUT.stem + f"_backup_final_{ts}.pptx")
    shutil.copy2(OUT, backup)
    prs.save(str(OUT))
    print(f"基于: {SRC.name}")
    print(f"备份: {backup.name}")
    print(f"已保存: {OUT} ({len(prs.slides)} 页)")


if __name__ == "__main__":
    main()
