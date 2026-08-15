#!/usr/bin/env python3
"""内容获客平台 PPT 全文核对：统一字体字号、优化文案、修复缺字体形状。"""

from __future__ import annotations

import re
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
DIR = ROOT / "docs/07-对外资料/03-汇报材料"
SRC = DIR / "内容获客平台_全域四平台实施方案_20260804.pptx"
BACKUP = DIR / "内容获客平台_全域四平台实施方案_20260804_backup_20260805.pptx"
OUT = SRC

FONT = "微软雅黑"
FOOTER = "内容获客平台 · AI 内容工厂 + 交易闭环 + CRM"

# 按形状 top 区间推断字号（EMU）
TITLE_TOP = (100000, 600000)      # 24pt
SUB_TOP = (450000, 700000)        # 12pt
SECTION_TOP = (900000, 5200000)   # 按宽度再分
FOOTER_TOP = (5900000, 6500000)   # 10pt

# 精确文案替换（全文匹配）
TEXT_REPLACEMENTS: dict[str, str] = {
    "虚拟/服务先行 · 实物门店后置\n决策层汇报 · 对标小鹅通/千聊 · 路径 A/B 可选":
        "虚拟与服务先行 · 实物与门店后置\n决策层汇报 · 对标小鹅通/千聊 · 路径 A/B 可选",
    "竞品调研 → 多品类方案 → 落地资质 → IP获客与闭环":
        "市场与竞品 → 三层方案 → 落地资质 → 获客闭环",
    "① 平台三层（AI 生产 + 交易 + CRM）  ② AI 内容智能体分期\n"
    "③ 商城 Phase1 首期范围  ④ 全域四平台挂载":
        "① 平台三层（AI 生产 + 交易 + CRM）\n"
        "② 智能体分期路线  ③ 商城首期范围  ④ 四平台分阶段挂载",
    "双链路 Mx → 政策 → 角色闭环 → 邀约 → 准备 → 技术 → 里程碑 → 差异化":
        "抖音 Mx → 政策资质 → 角色分工 → 邀约对照 → 准备清单 → 技术落点 → 里程碑",
    "IP内容矩阵 → Plan B → 场景 → 启动":
        "AI 内容生产 → IP 引流 → 风险预案 → 演示场景",
    "市场判断 · 为什么不做「另一个小鹅通」":
        "市场判断 · 为什么不做「另一个小鹅通」",
    "知识付费仍大 · 流量红利消退 · 差异化切入才有窗口":
        "知识付费市场仍在 · 流量红利见顶 · 差异化切入才有窗口",
    "轻量上手 + 多平台能挂 + 买完能履约 · 不拼功能大全":
        "轻量上手 · 多平台能挂 · 买完能履约 —— 不拼功能大全",
    "不正面打小鹅通广度 · 不硬拼千聊最低价 · 做内容获客的中间带":
        "不正面拼功能广度 · 不硬拼最低价 · 做「内容获客」中间带",
    "学履约闭环 · 不学功能大全":
        "学「订单→领权→履约」闭环 · 不学功能大全",
    "学官方店降门槛 · 警惕品牌与 CRM 短板":
        "学官方店降门槛 · 补足品牌与 CRM 短板",
    "店铺归属决定资质、品牌与数据 —— 两条路都要能讲清楚":
        "店铺归谁，决定资质、品牌与数据 —— 两条路径都要讲清楚",
    "内容获客平台（虚拟/服务）+ 分阶段四平台 · 课为首发 SKU":
        "虚拟/服务先行 · 课为首发 SKU · 四平台分阶段扩展",
    "公域能挂、店里能卖虚拟/服务、买完能履约；课为首发。商城装修/多端开店 Phase 2–3 补齐，再扩实物与门店 —— 路径 A/B 不写死。":
        "公域能挂、店里能卖、买完能履约。装修与多端开店 Phase 2–3 补齐，再扩实物与门店；路径 A/B 不写死。",
    "AI 生产层 · 公域层 · 内容获客平台 · 经营层（CRM）":
        "AI 内容工厂（生产）· 公域挂载 · 交易平台 · CRM 经营",
    "能生产获客内容 · 能挂四平台成交 · 能进 CRM 持续经营":
        "能生产内容 · 能挂平台成交 · 能进 CRM 持续经营",
    "战略一体、落地分期 —— 不与商城 Phase 1 抢工期":
        "战略一体、落地分期 —— 不与商城首期抢工期",
    "已有：多平台文案 / 笔记 / 视频脚本":
        "已交付：多平台文案、笔记、视频脚本",
    "智能体复用导出上架 · 不做配图/成片硬验收":
        "创作成果可导出复用 · 配图/成片不进硬验收",
    "课纲+详情+卖点一键生成商品 · 建议 v1.7":
        "课纲、详情、卖点一键生成商品草稿（建议 v1.7）",
    "做得快，更要用得稳 —— 降低公域发布侵权与平台限流风险":
        "做得快，更要发得稳 —— 降低侵权与平台限流风险",
    "仅商用授权白名单：图库 / 视频 / 音乐曲库":
        "仅用商用授权白名单：图库、视频、音乐曲库",
    "人审确认后导出 · 不承诺全自动发布":
        "人审确认后导出 · 不承诺全自动公域发布",
    "上架 → 付钱 → 履约（学/核销）→ 退款关权益":
        "上架 → 付款 → 履约（学习/核销）→ 退款关权益",
    "抖音优先 · 每平台先跑通买→履约→退 · 不齐发":
        "抖音优先 · 每平台先跑通「买→履约→退」· 不齐发",
    "产品同时支持 · 按客户资质与阶段选型":
        "两条路径并存 · 按客户资质与阶段选型",
    "两条官方技术公路 · 可替代抖店收银 ≠ 不用任何资质":
        "两条官方技术通道 · 可替代抖店收银，但仍需平台资质",
    "公域卖课多落教培强监管 · 其他虚拟/服务按类目合规":
        "公域卖课受教培强监管 · 其他虚拟/服务按类目合规",
    "路径 A · 链路 ① 抖店（默认演示包）· 客户轻 / 平台重":
        "路径 A + 链路 ① 抖店（默认演示）· 客户轻、平台重",
    "统一交易层 · 履约按 type 扩展 · 不污染 B2B 客户库":
        "统一交易层 · 履约按类型扩展 · 买家与 B2B 客户隔离",
    "支付闭环不可跳过 · 抖音公域 ①或② 至少一条验收":
        "支付闭环不可跳过 · 抖音公域至少验收一条链路",
    "相对千聊 / 小鹅通 / 有赞微盟":
        "相对千聊、小鹅通、有赞/微盟",
    "文案→图文→短视频 | 商用素材合规":
        "文案 → 图文 → 短视频\n商用素材可溯源",
    "获客打法 · AI 内容工厂 + IP 矩阵 → 内容获客平台":
        "获客打法 · AI 内容工厂 + IP 矩阵",
    "AI 辅助生产内容 → 系列视频引流 → 店内成交 → 自有端履约":
        "AI 辅助生产 → 系列视频引流 → 平台成交 → 自有端履约",
    "政策 · 同质化 · 摊大饼 —— 预先写好退路":
        "政策、同质化、范围失控 —— 预先写好退路",
    "多品类 + 品牌可升级 + AI内容工厂 + 商用素材合规":
        "多品类权益 · 品牌可升级 · AI 内容工厂 · 素材合规",
    "停四平台与实物扩张 · 收窄垂直 · 或做「已有小鹅通/千聊商家的经营增强层」验证 PMF。":
        "暂停四平台扩张 · 收窄垂直 · 或做「已有小鹅通/千聊商家的经营增强层」",
    "内容获客平台（虚拟/服务）+ 分阶段四平台 · IP获客 · A/B并列":
        "AI 内容工厂 + 交易闭环 + CRM · 四平台分阶段 · 路径 A/B 并列",
    "AI内容工厂+虚拟服务 | 交易闭环优先":
        "AI 内容工厂 + 虚拟/服务\n交易闭环优先",
    "AI生产→IP引流 | →内容获客平台":
        "AI 生产 → IP 引流\n→ 平台成交履约",
    "素材版权风险 | --- | 商用白名单曲库/图库 · 授权可追溯 · 人审发布 | 禁止无授权流行曲/爬图 | --- | Plan B（M6 未达标）":
        "素材版权：商用白名单 · 授权可追溯 · 人审发布\n禁止无授权流行曲与 AI 图商用",
    "① 平台三层（AI 生产 + 交易 + CRM） ② AI 内容智能体分期 | ③ 商城 Phase1 首期范围 ④ 全域四平台挂载":
        "① 平台三层  ② 智能体分期\n③ 商城首期范围  ④ 四平台分阶段挂载",
}

# 按 slide 索引（0-based）的局部替换
SLIDE_OVERRIDES: list[tuple[int, str, str]] = [
    (24, "与千聊同质化 / 素材侵权", "与千聊同质化"),
    (24, "素材版权风险 | --- | 商用白名单曲库/图库 · 授权可追溯 · 人审发布 | 禁止无授权流行曲/爬图 | --- | Plan B（M6 未达标）",
     "商用白名单图库/曲库 · 授权可追溯 · 人审后发布\n禁止无授权流行曲与 AI 图商用"),
]


def infer_size(shape, text: str, slide_idx: int) -> Pt:
    top = shape.top
    t = text.strip()
    if FOOTER in t:
        return Pt(10)
    # 页标题（内容页顶部大标题）
    if 100000 <= top <= 700000 and (
        "·" in t or t in {"内容框架", "总结与下一步"} or t.endswith("）") or len(t) <= 30
    ) and not t.startswith("①"):
        if slide_idx == 0 and "内容获客平台" in t:
            return Pt(28)
        return Pt(24)
    # 页副标题
    if 400000 <= top <= 900000 and len(t) < 80 and "分钟" not in t:
        if slide_idx > 0 and not re.fullmatch(r"0[1-4]", t):
            return Pt(12)
    if top >= FOOTER_TOP[0]:
        return Pt(10)
    if re.fullmatch(r"0[1-4]", t):
        return Pt(22)
    if t in {"M1", "M3", "M4–5", "Mx", "M6"}:
        return Pt(16)
    if t.startswith("Phase "):
        return Pt(16)
    if re.fullmatch(r"[\d,]+", t.replace("元", "").replace("+", "")) or t.endswith("元/年"):
        return Pt(20)
    if t in {"→", "›"}:
        return Pt(18)
    if slide_idx == 0 and "AI 内容工厂" in t:
        return Pt(14)
    if len(t) <= 10 and top > 4500000 and top < 5800000:
        return Pt(13)
    if len(t) <= 12 and top < 3500000:
        return Pt(13)
    return Pt(12)


def infer_bold(text: str, size: Pt) -> bool:
    t = text.strip()
    if size >= Pt(20):
        return True
    if size == Pt(24):
        return True
    if t.endswith("：") or t in {"优势 Strengths", "短板 Weaknesses"}:
        return True
    if re.match(r"^0[1-4]\s", t) or t.startswith("路径"):
        return True
    return False


def infer_color(text: str, size: Pt, slide_idx: int):
    t = text.strip()
    # 封面青色副标题
    if slide_idx == 0 and "AI 内容工厂" in t and size == Pt(14):
        return RGBColor(0x00, 0xC8, 0xFF)
    # 深色底上的浅色字（护城河 AI 列）
    if slide_idx == 22 and "文案" in t:
        return RGBColor(0xE2, 0xE8, 0xF0)
    if size <= Pt(10):
        return RGBColor(0x94, 0xA3, 0xB8)
    return None


def apply_run_style(run, size: Pt, bold: bool, color) -> None:
    run.font.name = FONT
    run.font.size = size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def set_shape_text(shape, new_text: str, slide_idx: int) -> None:
    if not shape.has_text_frame:
        if hasattr(shape, "text"):
            shape.text = new_text
        return
    size = infer_size(shape, new_text, slide_idx)
    bold = infer_bold(new_text, size)
    color = infer_color(new_text, size, slide_idx)
    tf = shape.text_frame
    tf.clear()
    lines = new_text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        for run in para.runs:
            apply_run_style(run, size, bold, color)


def normalize_shape(shape, slide_idx: int) -> None:
    if not shape.has_text_frame:
        return
    text = shape.text_frame.text
    if not text or not text.strip():
        return
    original = text.strip()
    new_text = TEXT_REPLACEMENTS.get(original, original)
    for idx, old, new in SLIDE_OVERRIDES:
        if slide_idx == idx and original == old:
            new_text = new
    # 统一页脚
    if "内容获客平台（虚拟/服务）+ 全域四平台实施方案" in new_text:
        new_text = FOOTER
    set_shape_text(shape, new_text, slide_idx)


def fix_slide9_subtitle(prs: Presentation) -> None:
    slide = prs.slides[8]
    for sh in slide.shapes:
        if sh.has_text_frame and sh.top < 700000 and "CRM" in sh.text_frame.text:
            set_shape_text(sh, "AI 内容工厂（生产）· 公域挂载 · 交易平台 · CRM 经营", 8)


def fix_slide23_moat(prs: Presentation) -> None:
    slide = prs.slides[22]
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t == "AI 内容工厂":
            set_shape_text(sh, "AI 内容工厂", 22)
        elif "文案" in t and "商用" in t:
            set_shape_text(sh, "文案 → 图文 → 短视频\n商用素材可溯源", 22)


def fix_slide24_title(prs: Presentation) -> None:
    slide = prs.slides[23]
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if "获客打法" in t:
            set_shape_text(sh, "获客打法 · AI 内容工厂 + IP 矩阵", 23)
        elif "AI 辅助" in t:
            set_shape_text(sh, "AI 辅助生产 → 系列视频引流 → 平台成交 → 自有端履约", 23)
        elif t.startswith("· 优先官方"):
            set_shape_text(
                sh,
                "· 优先官方挂载，少裸贴微信\n"
                "· IP 可多人入驻，勿绑死单一账号\n"
                "· Phase 2：AI 生成商品详情与推广素材\n"
                "· Phase 3/4：商用配图与短视频成片",
                23,
            )


def fix_slide1_cover(prs: Presentation) -> None:
    slide = prs.slides[0]
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t.startswith("内容获客平台"):
            set_shape_text(sh, "内容获客平台\n+ 全域四平台", 0)
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(28)
                    run.font.bold = True
        elif t == "AI 内容工厂 · 交易闭环 · CRM":
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(14)
                    run.font.color.rgb = RGBColor(0x00, 0xC8, 0xFF)


def fix_slide25_risk(prs: Presentation) -> None:
    slide = prs.slides[24]
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if "---" in t and ("素材" in t or "版权" in t):
            set_shape_text(
                sh,
                "素材版权\n商用白名单图库/曲库 · 授权可追溯 · 人审后发布\n禁止无授权流行曲与 AI 图商用",
                24,
            )
        elif t.startswith("暂停四平台") or t.startswith("停四平台"):
            set_shape_text(
                sh,
                "Plan B（M6 未达标）\n暂停四平台扩张 · 收窄垂直 · 或做「已有小鹅通/千聊商家的经营增强层」",
                24,
            )


def fix_slide27(prs: Presentation) -> None:
    slide = prs.slides[26]
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        top = sh.top
        if 450000 < top < 650000 and "交易闭环" in t and top < 600000:
            set_shape_text(
                sh,
                "AI 内容工厂 + 交易闭环 + CRM · 四平台分阶段 · 路径 A/B 并列",
                26,
            )
        elif top > 5800000:
            set_shape_text(sh, FOOTER, 26)
        elif t.startswith("① 锁定") or "M1 商户号" in t:
            set_shape_text(
                sh,
                "① 锁定首发品类（课 + 服务/数字权益）与路径 A/B\n"
                "② M1 商户号与资质；选定链路 ① 或 ②\n"
                "③ M3 支付验收 → Mx 公域首单\n"
                "④ Phase 2 智能体×商城 → Phase 3/4 图文与短视频",
                26,
            )
        elif top > 3500000 and top < 4500000 and "方案" in t or (
            "虚拟" in t and "交易" in t and len(t) < 40
        ):
            set_shape_text(sh, "AI 内容工厂 + 虚拟/服务\n交易闭环优先", 26)
        elif "IP 引流" in t or ("AI 生产" in t and "IP" in t):
            set_shape_text(sh, "AI 生产 → IP 引流\n→ 平台成交履约", 26)


def process_presentation(prs: Presentation) -> None:
    for i, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            normalize_shape(sh, i)
    fix_slide1_cover(prs)
    fix_slide9_subtitle(prs)
    fix_slide23_moat(prs)
    fix_slide24_title(prs)
    fix_slide25_risk(prs)
    fix_slide27(prs)


def save_prs(prs: Presentation) -> Path:
    for path in (OUT, DIR / "内容获客平台_全域四平台实施方案_20260804_无页码.pptx"):
        try:
            prs.save(str(path))
            print(f"已保存：{path}")
            return path
        except PermissionError:
            continue
    alt = DIR / "内容获客平台_全域四平台实施方案_20260804_final.pptx"
    prs.save(str(alt))
    print(f"主文件被占用，已另存：{alt}")
    return alt


def main() -> None:
    if not SRC.exists():
        raise SystemExit(f"缺少：{SRC}")
    # 备份
    stamp = DIR / "内容获客平台_全域四平台实施方案_20260804_pre_final.pptx"
    if not stamp.exists():
        shutil.copy2(SRC, stamp)

    prs = Presentation(str(SRC))
    process_presentation(prs)
    save_prs(prs)


if __name__ == "__main__":
    main()
