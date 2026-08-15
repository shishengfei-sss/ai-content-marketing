#!/usr/bin/env python3
"""移除 PPT 底部空黑块，统一异常页脚为青色细条。"""

from __future__ import annotations

import shutil
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu

ROOT = Path(__file__).resolve().parents[1]
PPT = ROOT / "docs/07-对外资料/03-汇报材料/内容获客平台_全域四平台实施方案_20260804.pptx"
FOOTER_TEXT = "内容获客平台 · AI 内容工厂 + 交易闭环 + CRM"
FOOTER_TOP = 6_537_960
FOOTER_HEIGHT = 27_432
CYAN = RGBColor(0x00, 0xC8, 0xFF)
DARK_NAVY = RGBColor(0x0A, 0x16, 0x28)
TEXT_DARK = RGBColor(0x11, 0x18, 0x27)


def remove_shape(shape) -> None:
    el = shape.element
    el.getparent().remove(el)


def is_empty_dark_pill(shape) -> bool:
    if not hasattr(shape, "fill"):
        return False
    try:
        if str(shape.fill.type) != "SOLID (1)":
            return False
        if str(shape.fill.fore_color.rgb) != "0A1628":
            return False
    except AttributeError:
        return False
    if shape.top < 5_800_000:
        return False
    text = shape.text.strip() if shape.has_text_frame else ""
    return not text


def is_bad_dark_footer(shape) -> bool:
    if not shape.has_text_frame or FOOTER_TEXT not in shape.text:
        return False
    if shape.height <= 100_000:
        return False
    try:
        if str(shape.fill.type) != "SOLID (1)":
            return True  # 非标准实心底栏也统一
        return str(shape.fill.fore_color.rgb) == "0A1628"
    except (AttributeError, TypeError):
        return True


def normalize_footer(shape) -> None:
    shape.top = FOOTER_TOP
    shape.height = FOOTER_HEIGHT
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = TEXT_DARK


def ensure_standard_footer(slide) -> None:
    footers = [
        s
        for s in slide.shapes
        if s.has_text_frame and FOOTER_TEXT in s.text
    ]
    if not footers:
        return
    primary = min(footers, key=lambda s: (abs(s.top - FOOTER_TOP), s.height))
    normalize_footer(primary)
    for s in footers:
        if s is not primary:
            remove_shape(s)


def fix_slide(slide) -> list[str]:
    actions: list[str] = []
    to_remove = []
    for sh in slide.shapes:
        if is_empty_dark_pill(sh):
            to_remove.append(sh)
        elif is_bad_dark_footer(sh):
            normalize_footer(sh)
            actions.append("页脚改为青色细条")
    for sh in to_remove:
        remove_shape(sh)
    if to_remove:
        actions.append(f"删除 {len(to_remove)} 个空黑占位块")
    ensure_standard_footer(slide)
    return actions


def main() -> None:
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    backup = PPT.with_name(PPT.stem + f"_backup_footerfix_{ts}.pptx")
    shutil.copy2(PPT, backup)
    prs = Presentation(str(PPT))
    for i, slide in enumerate(prs.slides):
        acts = fix_slide(slide)
        if acts:
            print(f"第 {i + 1} 页: {', '.join(acts)}")
    prs.save(str(PPT))
    print(f"备份: {backup.name}")
    print(f"已保存: {PPT}")


if __name__ == "__main__":
    main()
