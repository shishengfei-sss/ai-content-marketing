#!/usr/bin/env python3
"""生成「内容获客平台 · 一阶段进展汇报」13 页 PPT。"""

from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs/07-对外资料/03-汇报材料/内容获客平台-一阶段进展汇报-20260818.pptx"

# Brand colors
C_PRIMARY = RGBColor(0x1A, 0x56, 0xDB)
C_DARK = RGBColor(0x1E, 0x29, 0x3B)
C_MUTED = RGBColor(0x64, 0x74, 0x8B)
C_ACCENT = RGBColor(0x0E, 0xA5, 0xE9)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)


def set_run_font(run, *, size=18, bold=False, color=None, name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def add_title_bar(slide, title: str, subtitle: str | None = None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    set_run_font(r, size=28, bold=True, color=C_WHITE)
    tf.margin_left = Inches(0.55)
    tf.margin_top = Inches(0.18)
    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(12.2), Inches(0.45))
        tf2 = box.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        set_run_font(r2, size=14, color=C_MUTED)


def add_bullets(slide, items: list[str], *, top=1.75, left=0.65, width=12.0, size=18, spacing=8):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(spacing)
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = f"• {item}"
        set_run_font(r, size=size, color=C_DARK)


def add_table(slide, headers: list[str], rows: list[list[str]], *, top=1.75, col_widths: list[float] | None = None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.55), Inches(top), Inches(12.2), Inches(0.55 * n_rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                set_run_font(r, size=14, bold=True, color=C_WHITE)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_LIGHT_BG
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    set_run_font(r, size=13, color=C_DARK)
    return table_shape


def add_code_block(slide, text: str, *, top=1.75):
    box = slide.shapes.add_shape(1, Inches(0.55), Inches(top), Inches(12.2), Inches(1.55))
    box.fill.solid()
    box.fill.fore_color.rgb = C_LIGHT_BG
    box.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_run_font(r, size=16, bold=True, color=C_DARK, name="Consolas")


def add_footer(slide, page: int, total: int = 13):
    box = slide.shapes.add_textbox(Inches(11.6), Inches(7.15), Inches(1.2), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{page:02d} / {total}"
    set_run_font(r, size=11, color=C_MUTED)


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_WHITE
    bg.line.fill.background()
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return slide


def build_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    today = date.today().strftime("%Y-%m-%d")

    # Slide 1 - Cover
    s = blank_slide(prs)
    band = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(4.6))
    band.fill.solid()
    band.fill.fore_color.rgb = C_PRIMARY
    band.line.fill.background()
    tbox = s.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.5), Inches(2.2))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "内容获客平台 · 一阶段进展汇报"
    set_run_font(r, size=40, bold=True, color=C_WHITE)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "从「能写内容」到「能卖课、能履约、能管单」"
    set_run_font(r2, size=22, color=RGBColor(0xDB, 0xEA, 0xFE))
    meta = s.shapes.add_textbox(Inches(0.8), Inches(5.1), Inches(11.5), Inches(1.2))
    mtf = meta.text_frame
    for line in [f"汇报日期：{today}", "演示环境：本机测试 / 内网 192.168.20.201", "Phase 1 · 商城交易 MVP"]:
        p = mtf.paragraphs[0] if line == f"汇报日期：{today}" else mtf.add_paragraph()
        rr = p.add_run()
        rr.text = line
        set_run_font(rr, size=16, color=C_MUTED)
    add_footer(s, 1)

    # Slide 2
    s = blank_slide(prs)
    add_title_bar(s, "一句话定位")
    add_bullets(
        s,
        [
            "核心：AI 内容工厂 + 交易闭环 + 客户经营闭环",
            "智营已有：AI 写文案/笔记/脚本 + 线索、客户、商机跟到回款",
            "一阶段补齐：交易层 —— 商家上架虚拟课/资料/服务，买家付钱、学课、核销",
            "对标小鹅通/千聊：不只卖「店」，而是 生产 → 成交 → 持续经营 一条链",
        ],
        top=1.85,
        size=20,
    )
    add_footer(s, 2)

    # Slide 3
    s = blank_slide(prs)
    add_title_bar(s, "解决什么问题")
    add_table(
        s,
        ["现状痛点", "一阶段目标"],
        [
            ["内容在文档里，成交在别的平台", "内容和成交在同一产品体系"],
            ["公域买了课，私域没法履约", "公域订单可领权 → 自有端学课/核销"],
            ["平台管不了商家与商品合规", "入驻审核 + 商品上架闸 + 公域挂载闸"],
            ["卖课工具与客户跟进各管各的", "同一租户体系，后续可打通经营数据"],
        ],
        top=1.85,
        col_widths=[5.8, 6.4],
    )
    add_footer(s, 3)

    # Slide 4
    s = blank_slide(prs)
    add_title_bar(s, "一阶段交付什么", "平台 → 商家（三类主体）→ 店铺（可多店）→ 商品 → 订单 → 权益")
    add_bullets(
        s,
        [
            "三端齐全：平台运营端 P01–P12 · 商家 Web A01–A23 · 买家 H5/小程序 M01–M15",
            "三类商家主体：个人 / 个体户 / 企业入驻；支持多店铺",
            "三类商品：专栏/单课 · 数字资料包 · 服务（预约/次数卡/核销）",
            "交易闭环：下单 → 支付（Mock）→ 权益开通 → 学课/下载/核销",
            "退款关权益：退款后不可再履约",
            "合规：商品机审 + 平台人审；未过审不可上架",
            "公域预留：抖音 Mx 链路（Mock 档可演示领权）",
            "SaaS 套餐：P10 配置 + P11 人工开通（验收期）",
        ],
        top=1.95,
        size=17,
        spacing=6,
    )
    add_footer(s, 4)

    # Slide 5
    s = blank_slide(prs)
    add_title_bar(s, "一阶段不做什么", "硬验收：入驻 → 商品合规 → 私域可售 → 履约 → 退款关权益 →（Mock）公域领权")
    add_table(
        s,
        ["本期不做", "说明"],
        [
            ["商家自助购套餐", "线下对公 + 平台人工开通"],
            ["真微信 / 真抖店支付", "演示用 Mock，真机另签"],
            ["实物 / 门店 SKU", "Phase 2+"],
            ["H5 独立商城装修", "买家端以小程序/H5 履约为主"],
            ["AI 一键生成商品并上架", "Phase 2 打通创作顾问"],
            ["四平台齐发", "抖音优先，其余后移"],
        ],
        top=2.0,
        col_widths=[4.8, 7.4],
    )
    add_footer(s, 5)

    # Slide 6
    s = blank_slide(prs)
    add_title_bar(s, "三端角色一览")
    add_code_block(
        s,
        "平台运营：审入驻 · 审商品 · 套餐 · 结算 · 稽查\n"
        "商家后台：上架 · 订单 · 买家 · 权益 · 核销 · 公域映射\n"
        "买家端：进店 · 下单 · 已购 · 学课 · 领权 · 发票",
        top=1.85,
    )
    add_table(
        s,
        ["端", "谁用", "本期核心价值"],
        [
            ["平台端", "超管/运营", "入驻审核、商品审核、订阅台账、渠道配置、违规稽查"],
            ["商家端", "教培/咨询商家", "商品、订单、买家、权益、核销、公域映射、多店管理"],
            ["买家端", "C 端学员", "进店、下单、已购、学课、领权、发票"],
        ],
        top=3.65,
        col_widths=[1.6, 2.4, 8.2],
    )
    add_footer(s, 6)

    # Slide 7
    s = blank_slide(prs)
    add_title_bar(s, "合规两道闸")
    add_bullets(
        s,
        [
            "素材合规（AI 内容工厂，已有基础）：选源 → 成片 → 发布，审的是「营销内容」",
            "交易合规（一阶段新增，审的是「可售商品」）：",
            "　① 主体资质闸：入驻验证照与主体类型",
            "　② 商品上架闸：机审 + 平台人审，不过不能卖",
            "　③ 公域挂载闸：店内可卖 ≠ 可挂抖店/课程库",
            "　④ 事后稽查：举报下架、审计留痕",
        ],
        top=1.9,
        size=18,
    )
    add_footer(s, 7)

    # Slide 8
    s = blank_slide(prs)
    add_title_bar(s, "与现有智营产品的关系")
    add_table(
        s,
        ["阶段", "内容", "状态"],
        [
            ["Phase 0", "AI 创作顾问 + 线索跟到回款", "✅ 已交付"],
            ["Phase 1", "商城交易 MVP", "🚧 本期汇报重点"],
            ["Phase 2", "创作成果一键变商品草稿", "规划中"],
            ["Phase 3+", "AI 商用配图、短视频半自动", "远期"],
        ],
        top=1.9,
        col_widths=[2.0, 6.8, 3.4],
    )
    note = s.shapes.add_textbox(Inches(0.55), Inches(4.55), Inches(12.2), Inches(1.2))
    p = note.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = (
        "一阶段不是重做客户经营模块，是在智营上长出「卖课收钱」这一块，"
        "与现有线索、商机、报价、回款同体系演进。"
    )
    set_run_font(r, size=16, bold=True, color=C_PRIMARY)
    add_footer(s, 8)

    # Slide 9
    s = blank_slide(prs)
    add_title_bar(s, "现场演示", "商家卖课 → 买家学课 → 退款关权")
    add_bullets(
        s,
        [
            "环境（二选一）：本机 Web 5173 · 买家 H5 5174 · API 8003",
            "　　　　　　　或 内网 192.168.20.201:8088 / :8089",
            "平台超管：13800000000 / admin123456",
            "主商家：13900000099 / test123456",
            "验证码：1111 · 支付 Mock 点确认即成功",
            "演示商品：演示·IP获客实战课 · 演示·话术模板资料包 · 演示·1v1咨询次数卡",
        ],
        top=1.95,
        size=18,
    )
    add_footer(s, 9)

    # Slide 10 - demo script as table
    s = blank_slide(prs)
    add_title_bar(s, "演示剧本（约 15 分钟）")
    add_table(
        s,
        ["步", "角色", "动作", "看点"],
        [
            ["①", "平台", "商家租户列表", "多状态商家"],
            ["②", "平台", "审核中入驻详情", "平台管准入"],
            ["③", "商家", "看板 + 商品列表", "三类在售 + 草稿不可售"],
            ["④", "买家", "选课 → Mock 支付", "成交发生"],
            ["⑤", "买家", "已购 → 学课/播放", "履约发生"],
            ["⑥", "买家/商家", "已付单退款", "权益立即关闭"],
            ["⑦", "买家", "（可选）公域领权", "公域 → 自有端履约"],
            ["⑧", "平台/商家", "订阅台账 / 套餐", "SaaS 套餐模型"],
        ],
        top=1.75,
        col_widths=[0.7, 1.5, 3.8, 6.2],
    )
    add_footer(s, 10)

    # Slide 11
    s = blank_slide(prs)
    add_title_bar(s, "当前进展")
    add_table(
        s,
        ["指标", "说明"],
        [
            ["三端页面", "平台 12 类 + 商家 23 类 + 买家 15 类（PRD 规格）"],
            ["研发批次", "M0–M7 Mock 主路径 ✅；开箱演示可跑"],
            ["测试", "R1 API + 联测金标准 + UX 走查收口中"],
            ["演示数据", "4 家经营中商家，每家 36～50 条商品/订单量级"],
        ],
        top=1.9,
        col_widths=[2.8, 9.4],
    )
    note = s.shapes.add_textbox(Inches(0.55), Inches(4.55), Inches(12.2), Inches(0.8))
    p = note.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "私域交易闭环可演示、可验收；真机微信/抖店为商务并行线，不阻塞本期签收。"
    set_run_font(r, size=17, bold=True, color=C_PRIMARY)
    add_footer(s, 11)

    # Slide 12
    s = blank_slide(prs)
    add_title_bar(s, "已知限制与下一步")
    add_table(
        s,
        ["已知限制", "说明", "挡演示？"],
        [
            ["支付是 Mock", "点确认即成功", "否"],
            ["套餐线下开通", "无商家自助收银台", "否"],
            ["结算人工确认", "非银行自动打款", "否"],
            ["购物车/物流/实物", "范围外", "否"],
        ],
        top=1.75,
        col_widths=[3.2, 6.0, 3.0],
    )
    add_bullets(
        s,
        [
            "下一步 ① 签收一阶段 Mock 演示版 → 对内培训、对客户预演",
            "下一步 ② 启动真机验收（微信商户号 / 抖店）",
            "下一步 ③ Phase 2：自助购套餐 / 创作打通商品 / 扩平台",
        ],
        top=4.35,
        size=16,
        spacing=5,
    )
    add_footer(s, 12)

    # Slide 13
    s = blank_slide(prs)
    add_title_bar(s, "总结")
    add_bullets(
        s,
        [
            "一阶段把智营从「内容 + 客户经营」补成能卖虚拟课/服务并履约的完整链路",
            "平台、商家、买家三端可演示，核心路径 下单 → 学课 → 退款关权 已打通",
            "下一步：Mock 签收 → 真机试点 → Phase 2 内容与商城打通",
        ],
        top=2.2,
        size=22,
        spacing=14,
    )
    thanks = s.shapes.add_textbox(Inches(0.55), Inches(5.6), Inches(12.2), Inches(0.6))
    p = thanks.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "谢谢 · Q&A"
    set_run_font(r, size=24, bold=True, color=C_ACCENT)
    add_footer(s, 13)

    return prs


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_prs()
    prs.save(str(OUT))
    print(f"Wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
