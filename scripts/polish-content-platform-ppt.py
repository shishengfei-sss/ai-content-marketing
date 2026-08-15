#!/usr/bin/env python3
"""核对并美化内容获客平台 PPT：新增/修改页与原版 slide9 版式一致。"""

from __future__ import annotations

import re
import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
DIR = ROOT / "docs/07-对外资料/03-汇报材料"
SRC = DIR / "内容获客平台_全域四平台实施方案_20260804_v2.pptx"
BACKUP = DIR / "内容获客平台_全域四平台实施方案_20260804_backup_20260805.pptx"
OUT = DIR / "内容获客平台_全域四平台实施方案_20260804.pptx"
TOTAL = 27

# slide9 模板坐标（EMU）
ROW_TOPS = [1051560, 2103120, 3154680]
LABEL_POS = (1005840, None)  # top = row_top + 137160
RIGHT_POS = (4389120, None)
RIGHT2_OFFSET = 367680  # 1554480-1188720
BOTTOM_HDR_TOP = 4160520
COL_HDR_TOP = 4572000
COL_BODY_TOP = 5029200
COL_LEFTS = [777240, 4480560, 8183880]
COL_BODY_LEFTS = [822960, 4526280, 8229600]


def remove_shape(shape) -> None:
    el = shape.element
    el.getparent().remove(el)


def find_shape(slide, top: int, left: int, tol: int = 80000):
    for sh in slide.shapes:
        if abs(sh.top - top) < tol and abs(sh.left - left) < tol:
            return sh
    return None


def apply_text_like(shape, text: str, ref_shape) -> None:
    """写入文本并尽量继承参考 shape 的字号/颜色。"""
    if not shape.has_text_frame:
        shape.text = text
        return
    ref_run = None
    if ref_shape and ref_shape.has_text_frame and ref_shape.text_frame.paragraphs:
        p0 = ref_shape.text_frame.paragraphs[0]
        if p0.runs:
            ref_run = p0.runs[0]
    tf = shape.text_frame
    tf.clear()
    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        if ref_run:
            for run in para.runs:
                run.font.name = ref_run.font.name or "微软雅黑"
                if ref_run.font.size:
                    run.font.size = ref_run.font.size
                run.font.bold = ref_run.font.bold
                if ref_run.font.color.type is not None:
                    try:
                        run.font.color.rgb = ref_run.font.color.rgb
                    except AttributeError:
                        pass


def reset_body_from_template(target, template) -> None:
    """保留页眉页脚，正文区用模板形状重建。"""
    for sh in list(target.shapes):
        if 850000 < sh.top < 6400000:
            remove_shape(sh)
    for sh in template.shapes:
        if 850000 < sh.top < 6400000:
            newel = deepcopy(sh.element)
            target.shapes._spTree.insert_element_before(newel, "p:extLst")


def set_header(target, template, title: str, subtitle: str) -> None:
    t_title = find_shape(template, 164592, 640080)
    t_sub = find_shape(template, 502920, 640080)
    s_title = find_shape(target, 164592, 640080)
    s_sub = find_shape(target, 502920, 640080)
    if s_title and t_title:
        apply_text_like(s_title, title, t_title)
    if s_sub and t_sub:
        apply_text_like(s_sub, subtitle, t_sub)


def fill_rows(slide, template, rows: list[tuple[str, str, str]]) -> None:
    """rows: (左标签, 右第一行, 右第二行)"""
    t_label = find_shape(template, 1188720, 1005840)
    t_right1 = find_shape(template, 1188720, 4389120)
    t_right2 = find_shape(template, 1554480, 4389120)
    for i, (label, r1, r2) in enumerate(rows[:3]):
        rt = ROW_TOPS[i]
        ly = rt + 137160
        r2y = rt + 137160 + RIGHT2_OFFSET
        sl = find_shape(slide, ly, 1005840)
        sr1 = find_shape(slide, ly, 4389120)
        sr2 = find_shape(slide, r2y, 4389120)
        if sl:
            apply_text_like(sl, label, t_label)
        if sr1:
            apply_text_like(sr1, r1, t_right1)
        if sr2:
            apply_text_like(sr2, r2, t_right2)


def fill_bottom_columns(slide, template, header: str, columns: list[tuple[str, str]]) -> None:
    t_hdr = find_shape(template, BOTTOM_HDR_TOP, 640080)
    t_ch = find_shape(template, COL_HDR_TOP, COL_LEFTS[0])
    t_cb = find_shape(template, COL_BODY_TOP, COL_BODY_LEFTS[0])
    sh = find_shape(slide, BOTTOM_HDR_TOP, 640080)
    if sh and t_hdr:
        apply_text_like(sh, header, t_hdr)
    for i, (title, body) in enumerate(columns[:3]):
        ch = find_shape(slide, COL_HDR_TOP, COL_LEFTS[i])
        cb = find_shape(slide, COL_BODY_TOP, COL_BODY_LEFTS[i])
        if ch and t_ch:
            apply_text_like(ch, title, t_ch)
        if cb and t_cb:
            apply_text_like(cb, body, t_cb)


def renumber_footers(prs: Presentation) -> None:
    for i, slide in enumerate(prs.slides):
        n = f"{i + 1:02d} / {TOTAL}"
        for sh in slide.shapes:
            if hasattr(sh, "text") and re.fullmatch(r"\d{2}\s*/\s*\d+", sh.text.strip()):
                sh.text = n


def fix_slide1(prs: Presentation, backup: Presentation) -> None:
    slide = prs.slides[0]
    b = backup.slides[0]
    b_sub = find_shape(b, 3108960, 640080)  # cyan subtitle ref - actually 502920 in backup is cyan
    b_cyan = find_shape(b, 3108960, 640080) or find_shape(b, 502920, 640080)
    # backup: shape 3 cyan at 3108960? check - backup shape 3 is 3108960 virtual - shape index 3 top 3108960 is wrong
    b_cyan = None
    b_gray = None
    for sh in b.shapes:
        if hasattr(sh, "text"):
            if sh.text.strip() == "虚拟/服务先行 · 实物门店后置":
                b_cyan = sh
            if "实施方案" in sh.text:
                b_gray = sh

    for sh in slide.shapes:
        if not hasattr(sh, "text"):
            continue
        t = sh.text.strip()
        if t == "AI 内容工厂 · 交易闭环 · CRM" and b_cyan:
            apply_text_like(sh, "AI 内容工厂 · 交易闭环 · CRM", b_cyan)
        elif "决策层 27 页" in t and b_gray:
            apply_text_like(
                sh,
                "虚拟/服务先行 · 实物门店后置\n决策层 27 页 · 对标小鹅通/千聊 · 路径 A/B 可选",
                b_gray,
            )
        elif t == "三位一体闭环":
            ref = find_shape(b, 4846320, 5623560) or find_shape(b, 4846320, 5623560)
            for bs in b.shapes:
                if hasattr(bs, "text") and bs.text.strip() == "双路径接入":
                    ref = bs
                    break
            apply_text_like(sh, "三位一体闭环", ref)
        elif t.startswith("AI 生产内容"):
            ref = find_shape(b, 5138928, 5623560)
            for bs in b.shapes:
                if hasattr(bs, "text") and "官方店" in bs.text:
                    ref = bs
                    break
            apply_text_like(sh, "AI 生产内容 → 平台成交履约 → CRM 持续经营", ref)
        elif re.fullmatch(r"\d{2}\s*/\s*\d+", t):
            ref = find_shape(b, 6035040, 640080)
            apply_text_like(sh, f"01 / {TOTAL}", ref)


def fix_slide2(prs: Presentation) -> None:
    slide = prs.slides[1]
    ref = find_shape(slide, 1645920, 1691640)  # section 01 desc
    for sh in slide.shapes:
        if hasattr(sh, "text") and "平台三层" in sh.text:
            apply_text_like(
                sh,
                "① 平台三层（AI 生产 + 交易 + CRM）  ② AI 内容智能体分期\n"
                "③ 商城 Phase1 首期范围  ④ 全域四平台挂载",
                ref,
            )


def rebuild_styled_slide(prs: Presentation, idx: int, title: str, subtitle: str,
                         rows: list[tuple[str, str, str]],
                         bottom_header: str, bottom_cols: list[tuple[str, str]]) -> None:
    template = prs.slides[8]
    target = prs.slides[idx]
    reset_body_from_template(target, template)
    set_header(target, template, title, subtitle)
    fill_rows(target, template, rows)
    fill_bottom_columns(target, template, bottom_header, bottom_cols)


def main() -> None:
    if not SRC.exists():
        raise SystemExit(f"缺少源文件：{SRC}")
    shutil.copy2(SRC, SRC.with_suffix(".pptx.working"))
    prs = Presentation(str(SRC))
    backup = Presentation(str(BACKUP))

    fix_slide1(prs, backup)
    fix_slide2(prs)

    rebuild_styled_slide(
        prs, 9,
        "平台三层 · AI 内容工厂 + 交易闭环 + CRM",
        "能生产获客内容 · 能挂四平台成交 · 能进 CRM 持续经营",
        [
            ("AI 内容工厂", "生产层 · 多平台文案 / 笔记 / 视频脚本", "知识库 + 品牌语气 + 合规自检"),
            ("内容获客平台", "交易层 · 商品 / 订单 / 支付 / 权益 · 四平台挂载", "Phase 1：能挂 · 能卖 · 能履约"),
            ("智营 CRM", "运营层（可选）· 活动 / 线索 / 客户跟进", "成交后经营 · 复用 v0.6～v1.6 营客能力"),
        ],
        "核心差异化",
        [
            ("竞品", "小鹅通/千聊：卖「店」与履约"),
            ("我们", "AI 生产内容 + 平台成交"),
            ("闭环", "CRM 持续经营"),
        ],
    )

    rebuild_styled_slide(
        prs, 10,
        "内容营销智能体 · 分期能力路线",
        "战略一体、落地分期 —— 不与商城 Phase 1 抢工期",
        [
            ("Phase 0 ✅", "已有：多平台文案 / 笔记 / 视频脚本", "v0.6 顾问 + v1.1 创作增强 · 可演示"),
            ("Phase 1", "商城 MVP：交易闭环优先", "智能体复用导出上架 · 不做配图/成片硬验收"),
            ("Phase 2", "智能体 × 商城打通", "课纲+详情+卖点一键生成商品 · 建议 v1.7"),
        ],
        "后续能力（建议版本）",
        [
            ("Phase 3", "AI 商用配图\n图文成品 ZIP\n建议 v1.8"),
            ("Phase 4", "短视频半自动\n商用 BGM + 人审导出\n建议 v1.9"),
            ("Phase 5", "效果闭环\nA/B · 归因 · Multi-Agent\n建议 v2.0+"),
        ],
    )

    rebuild_styled_slide(
        prs, 11,
        "素材合规 · 图 / 视 / 音三道闸",
        "做得快，更要用得稳 —— 降低公域发布侵权与平台限流风险",
        [
            ("选源闸", "仅商用授权白名单：图库 / 视频 / 音乐曲库", "AI 生成走商用套餐 · 上传须声明授权"),
            ("成片闸", "每份素材带 license 元数据", "未授权素材不可进入导出包"),
            ("发布闸", "文案合规 + 素材合规", "人审确认后导出 · 不承诺全自动发布"),
        ],
        "分阶段开放",
        [
            ("Phase 3", "商用配图\n多图 ZIP"),
            ("Phase 4", "商用 BGM\n半自动成片"),
            ("禁止项", "流行曲随意拼\n无授权 AI 图商用"),
        ],
    )

    # 同步 slide9 副标题（若尚未更新）
    s9 = prs.slides[8]
    sub = find_shape(s9, 502920, 640080)
    if sub and "AI 生产层" not in sub.text:
        t = find_shape(s9, 502920, 640080)
        apply_text_like(sub, "AI 生产层 · 公域层 · 内容获客平台 · 经营层（CRM）", t)

    renumber_footers(prs)

    saved = False
    for path in (OUT, SRC):
        try:
            prs.save(str(path))
            print(f"已保存：{path}")
            saved = True
            break
        except PermissionError:
            continue
    if not saved:
        alt = DIR / "内容获客平台_全域四平台实施方案_20260804_polished.pptx"
        prs.save(str(alt))
        print(f"文件被占用，已另存：{alt}")


if __name__ == "__main__":
    main()
