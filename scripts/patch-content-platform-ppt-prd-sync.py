#!/usr/bin/env python3
"""安全补丁：同步 PRD 架构（交易合规闸、商家多店套餐、主体类型）到决策层 PPT。仅文本替换，不 clear 版式。"""

from __future__ import annotations

import shutil
from datetime import datetime
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
PPT = ROOT / "docs/07-对外资料/03-汇报材料/内容获客平台_全域四平台实施方案_20260804.pptx"


def replace_in_slide(slide, mapping: dict[str, str]) -> int:
    n = 0
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = shape.text
        new = text
        for old, val in mapping.items():
            if old in new:
                new = new.replace(old, val)
        if new != text:
            shape.text = new
            n += 1
    return n


def main() -> None:
    if not PPT.exists():
        raise SystemExit(f"找不到 {PPT}")
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    backup = PPT.with_name(PPT.stem + f"_backup_prd_sync_{ts}.pptx")
    shutil.copy2(PPT, backup)
    prs = Presentation(str(PPT))
    total = 0

    # Slide 10: 三层 · 交易层
    total += replace_in_slide(
        prs.slides[9],
        {
            "交易层 · 商品 / 订单 / 支付 / 权益 · 四平台挂载": (
                "交易层 · 商家入驻/多店/套餐 · 商品/订单/支付/权益"
            ),
            "Phase 1：能挂 · 能卖 · 能履约": "Phase 1：能审 · 能卖 · 能履约 · 套餐可灵活配置",
        },
    )

    # Slide 13: Phase1 范围
    total += replace_in_slide(
        prs.slides[12],
        {
            "· 公域订单回流（抖音优先）": (
                "· 公域订单回流（抖音优先）\n"
                "· 商品上架闸：机审 + 平台人审\n"
                "· 公域挂载闸：过审后才可挂抖店/课程库"
            ),
            "硬验收：能上架虚拟/服务 → 能付钱 → 能履约。商城装修与多端开店 Phase 1 可不做，Phase 2/3 必补。": (
                "硬验收：商家入驻 → 商品合规(机审+人审) → 私域可售 → 公域挂载闸 → "
                "付钱履约 → 退款关权益。商城装修 Phase 2 补。"
            ),
        },
    )

    # Slide 18: 角色闭环
    total += replace_in_slide(
        prs.slides[17],
        {
            "① 在我们平台入驻": "① 入驻（个人/个体/企业）并开店",
            "③ 等内容审核通过": "③ 商品合规审核通过（机审+人审）",
            "+ 讲师入驻与课程审核": "+ 商家入驻 · 商品合规 · 套餐权益管理",
        },
    )

    # Slide 20: 平台准备
    total += replace_in_slide(
        prs.slides[19],
        {
            "· 讲师入驻 + 课程合规审核": "· 商家入驻(三类主体) + 商品合规 + 套餐配置",
            "· 种子讲师入驻 SOP（入驻→上架→挂车）": "· 种子商家 SOP（入驻→开店→过审上架→挂车）",
        },
    )

    # Slide 23: 差异化
    total += replace_in_slide(
        prs.slides[22],
        {
            "商用素材合规": "商用素材合规\nSaaS 套餐可配",
        },
    )

    # Slide 27: 总结
    total += replace_in_slide(
        prs.slides[26],
        {
            "AI 内容工厂 + 虚拟/服务\n交易闭环优先": "商家多店 + 套餐 SaaS\n交易合规 + 虚拟/服务闭环",
        },
    )

    prs.save(str(PPT))
    print(f"备份: {backup}")
    print(f"已替换 shape 数: {total}")
    print(f"已保存: {PPT}")


if __name__ == "__main__":
    main()
