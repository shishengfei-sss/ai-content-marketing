#!/usr/bin/env python3
"""将 PRD 架构（交易合规、商家多店、套餐）融入决策层 PPT，保持原稿版式与字体。"""

from __future__ import annotations

import shutil
import sys
from copy import deepcopy
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
DIR = ROOT / "docs/07-对外资料/03-汇报材料"
PPT = DIR / "内容获客平台_全域四平台实施方案_20260804.pptx"

FONT = "微软雅黑"
FOOTER = "内容获客平台 · AI 内容工厂 + 交易闭环 + CRM"
TOTAL_AFTER = 28


def duplicate_slide(prs: Presentation, source_index: int, insert_index: int):
    source = prs.slides[source_index]
    dest = prs.slides.add_slide(source.slide_layout)
    for shape in source.shapes:
        newel = deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(newel, "p:extLst")
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    new_id = slides[-1]
    xml_slides.remove(new_id)
    xml_slides.insert(insert_index, new_id)
    return prs.slides[insert_index]


def get_shape_by_pos(slide, top_min: int, top_max: int, left_max: int = 2_000_000):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if top_min <= sh.top <= top_max and sh.left < left_max:
            return sh
    return None


def get_shape_right(slide, top_min: int, top_max: int, left_min: int = 4_000_000):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if top_min <= sh.top <= top_max and sh.left >= left_min:
            return sh
    return None


def replace_preserve(shape, new_text: str) -> None:
    if not shape or not shape.has_text_frame:
        return
    ref = None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            ref = run
            break
        if ref:
            break
    size = ref.font.size if ref and ref.font.size else Pt(12)
    bold = ref.font.bold if ref else False
    color = RGBColor(0x11, 0x18, 0x27)
    if ref:
        try:
            color = ref.font.color.rgb
        except AttributeError:
            pass
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(new_text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.name = FONT
            run.font.size = size
            run.font.bold = bold
            try:
                run.font.color.rgb = color
            except AttributeError:
                pass


def set_title_subtitle(slide, title: str, subtitle: str) -> None:
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if 140_000 < sh.top < 220_000 and sh.left < 1_500_000:
            replace_preserve(sh, title)
        elif 450_000 < sh.top < 650_000 and sh.left < 1_500_000:
            replace_preserve(sh, subtitle)


def fill_gate_slide(slide) -> None:
    """沿用素材合规页版式：三行主栏 + 底栏三列（第四道闸并入第三行）。"""
    set_title_subtitle(
        slide,
        "交易合规 · 可售商品四道闸",
        "卖什么要能审 —— 与「素材合规」分工：素材管创作导出，交易管可售 SKU",
    )
    rows = [
        (
            1_100_000,
            1_700_000,
            "主体资质闸",
            "个人 / 个体工商户 / 企业 分类入驻\n执照或身份证 · 支付进件",
        ),
        (
            2_150_000,
            2_750_000,
            "商品上架闸",
            "机审 + 平台人审 · 通过后方可私域上架\n未过审不可上架销售",
        ),
        (
            3_200_000,
            3_900_000,
            "公域与稽查",
            "公域挂载闸：过上架闸后才可同步抖店 / 课程库\n"
            "事后稽查闸：举报抽检 · 强制下架（存量订单照常履约）",
        ),
    ]
    for top_min, top_max, label, body in rows:
        replace_preserve(get_shape_by_pos(slide, top_min, top_max), label)
        rights = sorted(
            [
                s
                for s in slide.shapes
                if s.has_text_frame and top_min <= s.top <= top_max and s.left >= 4_000_000
            ],
            key=lambda s: s.top,
        )
        if rights:
            replace_preserve(rights[0], body)
            for extra in rights[1:]:
                replace_preserve(extra, "")

    for sh in slide.shapes:
        if sh.has_text_frame and sh.text.strip() == "分阶段开放":
            replace_preserve(sh, "平台能力（与素材闸并列）")
            break

    cols = sorted(
        [s for s in slide.shapes if s.has_text_frame and s.top > 4_500_000 and s.width > 2_500_000],
        key=lambda s: s.left,
    )
    titles = sorted([c for c in cols if c.height < 400_000], key=lambda s: s.left)
    bodies = sorted([c for c in cols if c.height >= 400_000], key=lambda s: s.left)
    bottom = [
        ("商家 · 多店", "入驻 → 多店铺 → 商品挂店下\n配额受当前套餐限制"),
        ("SaaS 套餐", "免费 / 基础 / 旗舰 平台可配\n按生效～失效周期管理权益"),
        ("与素材闸", "素材：审创作导出物\n交易：审可售卖商品"),
    ]
    for i, (lab, bod) in enumerate(bottom):
        if i < len(titles):
            replace_preserve(titles[i], lab)
        if i < len(bodies):
            replace_preserve(bodies[i], bod)


def replace_in_slide(slide, mapping: dict[str, str]) -> None:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text
        new = text
        for old, val in mapping.items():
            if old in new:
                new = new.replace(old, val)
        if new != text:
            replace_preserve(shape, new)


def has_trade_compliance_slide(prs: Presentation) -> bool:
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame and "可售商品四道闸" in sh.text:
                return True
    return False


def find_slide_by_title(prs: Presentation, keyword: str) -> int | None:
    for i, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if sh.has_text_frame and keyword in sh.text and sh.top < 300_000:
                return i
    return None


def polish_slides(prs: Presentation) -> None:
    """轻量润色：与新增页语气一致，避免生硬堆术语。"""
    i_dir = find_slide_by_title(prs, "内容框架") or 1
    i_arch = find_slide_by_title(prs, "系统方案总览") or 8
    i_dual = find_slide_by_title(prs, "平台三层") or 9
    i_phase1 = find_slide_by_title(prs, "Phase 1 范围") or 12
    i_role = find_slide_by_title(prs, "角色分工") or 17
    i_prep = find_slide_by_title(prs, "平台方准备清单") or 19
    i_diff = find_slide_by_title(prs, "差异化与护城河") or 22
    i_sum = find_slide_by_title(prs, "总结与下一步") or 26

    replace_in_slide(
        prs.slides[i_dir],
        {
            "① 平台三层（AI 生产 + 交易 + CRM） ② AI 内容智能体分期\n"
            "③ 商城 Phase1 首期范围  ④ 全域四平台挂载": (
                "① 平台三层（商家多店）  ② 智能体分期\n"
                "③ 商城首期 + 交易合规  ④ 四平台分阶段挂载"
            ),
            "① 平台三层  ② 智能体分期\n③ 商城首期范围  ④ 四平台分阶段挂载": (
                "① 平台三层（商家多店）  ② 智能体分期\n"
                "③ 商城首期 + 交易合规  ④ 四平台分阶段挂载"
            ),
        },
    )

    replace_in_slide(
        prs.slides[i_arch],
        {
            "统一商品/订单/支付/权益 · 履约按类型插件": (
                "商家入驻 · 多店铺 · 套餐 SaaS · 统一交易层"
            ),
            "课·数字权益·预约核销（初期）": "商品/订单/支付/权益 · 课·数字·服务履约",
        },
    )

    replace_in_slide(
        prs.slides[i_dual],
        {
            "Phase 1：能审 · 能卖 · 能履约 · 套餐可灵活配置": (
                "Phase 1：入驻过审 · 商品合规 · 能卖能履约"
            ),
            "交易层 · 商家入驻/多店/套餐 · 商品/订单/支付/权益": (
                "交易层 · 商家多店 · 套餐 SaaS · 商品/订单/权益"
            ),
            "我们\nAI 生产内容 + 平台成交": "我们\nAI 生产 + 多店成交 + 套餐可扩展",
            "AI 生产内容 + 平台成交": "AI 生产 + 多店成交 + 套餐可扩展",
        },
    )

    replace_in_slide(
        prs.slides[i_phase1],
        {
            "上架 → 付钱 → 履约（学/核销）→ 退款关权益": (
                "入驻 → 合规上架 → 付钱 → 履约 → 退款关权益"
            ),
            "· 公域订单回流（抖音优先）": (
                "· 公域订单回流（抖音优先）\n"
                "· 商品合规：机审 + 平台人审\n"
                "· 公域挂载：过审后才可挂抖店"
            ),
            "硬验收：商家入驻 → 商品合规(机审+人审) → 私域可售 → 公域挂载闸 → 付钱履约 → 退款关权益。商城装修 Phase 2 补。": (
                "硬验收：入驻 → 商品合规 → 私域可售 → 公域挂载 → 履约 → 退款关权益。"
                "装修开店 Phase 2 补。"
            ),
            "· 店铺装修 / 开店向导 → P2": "· 多店铺 · 套餐自助购 → P2\n· 店铺装修 / 向导 → P2",
        },
    )

    replace_in_slide(
        prs.slides[i_role],
        {
            "① 在我们平台入驻": "① 平台入驻（个人/个体/企业）",
            "① 入驻（个人/个体/企业）并开店": "① 平台入驻（个人/个体/企业）",
            "② 建课 / 服务 / 定价": "② 开店 · 建课/服务 · 定价",
            "③ 等内容审核通过": "③ 商品合规审核通过",
            "③ 商品合规审核通过（机审+人审）": "③ 商品合规审核通过",
        },
    )

    replace_in_slide(
        prs.slides[i_prep],
        {
            "· 讲师入驻 + 课程合规审核": "· 商家入驻 + 商品合规 + 套餐配置",
            "· 种子讲师入驻 SOP（入驻→上架→挂车）": "· 种子商家 SOP（入驻→开店→合规上架→挂车）",
            "· 种子商家 SOP（入驻→开店→过审上架→挂车）": "· 种子商家 SOP（入驻→开店→合规上架→挂车）",
        },
    )

    replace_in_slide(
        prs.slides[i_diff],
        {
            "商用素材合规": "商用素材合规\n套餐 SaaS 可配",
            "· 不拼四平台最低价": "· 不拼四平台最低价\n· 套餐分层（对标小鹅通年费档）",
        },
    )

    replace_in_slide(
        prs.slides[i_sum],
        {
            "AI 内容工厂 + 虚拟/服务\n交易闭环优先": (
                "商家多店 · 套餐 SaaS\n交易合规 + 闭环优先"
            ),
            "商家多店 + 套餐 SaaS\n交易合规 + 虚拟/服务闭环": (
                "商家多店 · 套餐 SaaS\n交易合规 + 闭环优先"
            ),
        },
    )


def fix_footers(slide) -> None:
    footers = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if not t:
            continue
        if FOOTER in t or (sh.top > 6_000_000 and len(t) < 80):
            footers.append(sh)
    if not footers:
        return
    primary = min(footers, key=lambda s: (s.top, s.left))
    replace_preserve(primary, FOOTER)
    for sh in footers[1:]:
        el = sh.element
        el.getparent().remove(el)


def audit(prs: Presentation) -> list[str]:
    issues = []
    for i, slide in enumerate(prs.slides):
        footers = 0
        big_bottom = []
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text.strip()
            if FOOTER in t:
                footers += 1
            if i > 0 and sh.top > 5_500_000 and "全域四平台" in t and FOOTER not in t:
                big_bottom.append(t[:40])
        if footers > 1:
            issues.append(f"S{i+1} 重复页脚({footers})")
        if big_bottom:
            issues.append(f"S{i+1} 底部异常标题: {big_bottom}")
    return issues


def extract_preview(prs: Presentation, path: Path) -> None:
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"=== Slide {i} ===")
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text.strip():
                lines.append(sh.text.strip())
        lines.append("")
    path.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    if not PPT.exists():
        raise SystemExit(f"找不到 {PPT}")

    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    backup = PPT.with_name(PPT.stem + f"_backup_integrate_{ts}.pptx")
    shutil.copy2(PPT, backup)

    prs = Presentation(str(PPT))

    has_trade = has_trade_compliance_slide(prs)
    if not has_trade:
        new_slide = duplicate_slide(prs, 11, 12)
        fill_gate_slide(new_slide)
    else:
        for slide in prs.slides:
            for sh in slide.shapes:
                if sh.has_text_frame and "可售商品四道闸" in sh.text:
                    fill_gate_slide(slide)
                    break

    polish_slides(prs)

    for slide in prs.slides:
        fix_footers(slide)

    issues = audit(prs)
    preview = ROOT / ".cursor" / "ppt-after-integrate.txt"
    extract_preview(prs, preview)

    try:
        prs.save(str(PPT))
        out = PPT
    except PermissionError:
        out = PPT.with_name(PPT.stem + "_integrated.pptx")
        prs.save(str(out))

    print(f"备份: {backup}")
    print(f"保存: {out}")
    print(f"总页数: {len(prs.slides)}")
    print(f"预览: {preview}")
    if issues:
        print("自检:")
        for x in issues:
            print(" ", x)
    else:
        print("自检通过")


if __name__ == "__main__":
    main()
