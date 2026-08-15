#!/usr/bin/env python3
"""修复内容获客 PPT：从 pre_final 底稿恢复样式，安全更新文案，逐页排查。"""

from __future__ import annotations

import re
import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
DIR = ROOT / "docs/07-对外资料/03-汇报材料"
BASE = DIR / "内容获客平台_全域四平台实施方案_20260804_pre_final.pptx"
BACKUP = DIR / "内容获客平台_全域四平台实施方案_20260804_backup_20260805.pptx"
OUT = DIR / "内容获客平台_全域四平台实施方案_20260804.pptx"

FONT = "微软雅黑"
FOOTER = "内容获客平台 · AI 内容工厂 + 交易闭环 + CRM"
OLD_FOOTER_MARKERS = (
    "内容获客平台（虚拟/服务）",
    "全域四平台实施方案",
)


def is_footer_text(text: str) -> bool:
    t = text.strip()
    if t == FOOTER:
        return True
    return any(m in t for m in OLD_FOOTER_MARKERS)

# 封面固定样式
COVER = {
    "title": ("内容获客平台\n+ 全域四平台", Pt(28), RGBColor(0xFF, 0xFF, 0xFF), True),
    "tagline": ("AI 内容工厂 · 交易闭环 · CRM", Pt(14), RGBColor(0x00, 0xC8, 0xFF), True),
    "desc": (
        "虚拟与服务先行 · 实物与门店后置\n决策层汇报 · 对标小鹅通/千聊 · 路径 A/B 可选",
        Pt(12),
        RGBColor(0xB0, 0xC4, 0xD8),
        False,
    ),
    "footer": (FOOTER, Pt(10), RGBColor(0x6B, 0x72, 0x80), False),
}
COVER_CARD_TITLE = (Pt(16), RGBColor(0x0A, 0x16, 0x28), True)
COVER_CARD_BODY = (Pt(12), RGBColor(0x6B, 0x72, 0x80), False)

# 全文精确替换（仅改字，样式由上下文恢复）
TEXT_MAP: dict[str, str] = {
    "虚拟/服务先行 · 实物门店后置\n决策层汇报 · 对标小鹅通/千聊 · 路径 A/B 可选":
        COVER["desc"][0],
    "内容获客平台（虚拟/服务）+ 全域四平台实施方案": FOOTER,
    "竞品调研 → 多品类方案 → 落地资质 → IP获客与闭环":
        "市场与竞品 → 三层方案 → 落地资质 → 获客闭环",
    "① 平台三层（AI 生产 + 交易 + CRM）  ② AI 内容智能体分期\n"
    "③ 商城 Phase1 首期范围  ④ 全域四平台挂载":
        "① 平台三层  ② 智能体分期\n③ 商城首期范围  ④ 四平台分阶段挂载",
    "AI 生产层 · 公域层 · 内容获客平台 · 经营层（CRM）":
        "AI 内容工厂（生产）· 公域挂载 · 交易平台 · CRM 经营",
    "能生产获客内容 · 能挂四平台成交 · 能进 CRM 持续经营":
        "能生产内容 · 能挂平台成交 · 能进 CRM 持续经营",
    "战略一体、落地分期 —— 不与商城 Phase 1 抢工期":
        "战略一体、落地分期 —— 不与商城首期抢工期",
    "已有：多平台文案 / 笔记 / 视频脚本":
        "已交付：多平台文案、笔记、视频脚本",
    "智能体复用导出上架 · 不做配图/成片硬验收":
        "创作成果可导出复用 · 配图/成片不进硬验收",
    "课纲+详情+卖点一键生成商品 · 建议 v1.7":
        "课纲、详情、卖点一键生成商品草稿（建议 v1.7）",
    "做得快，更要用得稳 —— 降低公域发布侵权与平台限流风险":
        "做得快，更要发得稳 —— 降低侵权与平台限流风险",
    "内容获客平台（虚拟/服务）+ 分阶段四平台 · IP获客 · A/B并列":
        "AI 内容工厂 + 交易闭环 + CRM · 四平台分阶段 · 路径 A/B 并列",
    "AI内容工厂+虚拟服务 | 交易闭环优先":
        "AI 内容工厂 + 虚拟/服务\n交易闭环优先",
    "AI生产→IP引流 | →内容获客平台":
        "AI 生产 → IP 引流\n→ 平台成交履约",
}


def remove_shape(shape) -> None:
    el = shape.element
    el.getparent().remove(el)


def write_styled(shape, text: str, size: Pt, color: RGBColor, bold: bool = False) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.name = FONT
            run.font.size = size
            run.font.bold = bold
            run.font.color.rgb = color


def replace_preserve(shape, new_text: str) -> None:
    """改字但尽量保留原 run 样式。"""
    if not shape.has_text_frame:
        return
    ref = None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            ref = run
            break
        if ref:
            break
    size = ref.font.size if ref and ref.font.size else Pt(12)
    bold = ref.font.bold if ref else False
    color = None
    if ref:
        try:
            color = ref.font.color.rgb
        except AttributeError:
            pass
    if color is None:
        color = RGBColor(0x11, 0x18, 0x27)
    write_styled(shape, new_text, size, color, bold)


def map_text(text: str) -> str:
    t = text.strip()
    if t in TEXT_MAP:
        return TEXT_MAP[t]
    for old, new in TEXT_MAP.items():
        if old in t:
            t = t.replace(old, new)
    if any(m in t for m in OLD_FOOTER_MARKERS):
        return FOOTER
    return t


COVER_CARD_TOPS = {1554480, 2651760, 3749040, 4846320}


def fix_cover(slide) -> None:
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        top, left = sh.top, sh.left
        if top < 3500000 and t.startswith("内容获客平台") and "+ 全域" in t:
            write_styled(sh, *COVER["title"])
            continue
        if "AI 内容工厂" in t and top < 3500000 and left < 2000000:
            write_styled(sh, *COVER["tagline"])
            continue
        if top < 4500000 and left < 2000000 and ("虚拟" in t or "决策层" in t):
            write_styled(sh, *COVER["desc"])
            continue
        if top > 5500000:
            write_styled(sh, *COVER["footer"])
            continue
        if left > 5000000 and top > 1000000:
            is_title = any(abs(top - tt) < 120000 for tt in COVER_CARD_TOPS)
            if is_title:
                write_styled(sh, t, *COVER_CARD_TITLE)
            else:
                write_styled(sh, t, *COVER_CARD_BODY)


def fix_footers(slide, slide_idx: int) -> None:
    """统一页脚，删除重复页脚形状。"""
    footers = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if not t:
            continue
        if is_footer_text(t):
            footers.append(sh)
        elif slide_idx > 0 and sh.top > 6000000 and sh.width > 5000000 and len(t) < 60:
            footers.append(sh)
        elif slide_idx == 0 and sh.top > 5800000 and is_footer_text(t):
            footers.append(sh)
    if not footers:
        return
    primary = min(footers, key=lambda s: (s.top, s.left))
    write_styled(primary, *COVER["footer"])
    for sh in footers:
        if sh is not primary:
            remove_shape(sh)


def fix_content_titles(slide) -> None:
    """内容页：顶部标题 24pt 白/深色，副标题 12pt。"""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if not t or is_footer_text(t):
            continue
        if 140000 < sh.top < 280000 and sh.left < 1500000 and len(t) < 50:
            # 深色顶栏上的白字标题
            write_styled(sh, map_text(t), Pt(24), RGBColor(0xFF, 0xFF, 0xFF), True)
        elif 450000 < sh.top < 650000 and sh.left < 1500000 and len(t) < 90:
            write_styled(sh, map_text(t), Pt(12), RGBColor(0x38, 0xBD, 0xF8), False)


def apply_text_map(slide) -> None:
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        raw = sh.text_frame.text.strip()
        if not raw:
            continue
        new = map_text(raw)
        if new != raw:
            replace_preserve(sh, new)


def fix_slide27(slide) -> None:
    apply_text_map(slide)
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t.startswith("① 锁定") or "M1 商户号" in t:
            replace_preserve(
                sh,
                "① 锁定首发品类（课 + 服务/数字权益）与路径 A/B\n"
                "② M1 商户号与资质；选定链路 ① 或 ②\n"
                "③ M3 支付验收 → Mx 公域首单\n"
                "④ Phase 2 智能体×商城 → Phase 3/4 图文与短视频",
            )
        elif "AI内容工厂" in t.replace(" ", "") and "虚拟" in t:
            replace_preserve(sh, "AI 内容工厂 + 虚拟/服务\n交易闭环优先")
        elif "AI生产" in t.replace(" ", "") or ("IP 引流" in t and "平台" in t):
            replace_preserve(sh, "AI 生产 → IP 引流\n→ 平台成交履约")


def fix_slide25(slide) -> None:
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if "---" in t and ("素材" in t or "版权" in t):
            replace_preserve(
                sh,
                "素材版权\n商用白名单图库/曲库 · 授权可追溯 · 人审后发布\n禁止无授权流行曲与 AI 图商用",
            )
        elif t.startswith("Plan B") or t.startswith("暂停") or t.startswith("停四平台"):
            replace_preserve(
                sh,
                "Plan B（M6 未达标）\n暂停四平台扩张 · 收窄垂直 · 或做「已有小鹅通/千聊商家的经营增强层」",
            )


def audit(slide, idx: int) -> list[str]:
    issues = []
    titles_bottom = []
    footers = []
    for j, sh in enumerate(slide.shapes):
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if not t:
            continue
        sizes = []
        colors = []
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    sizes.append(run.font.size.pt)
                try:
                    if run.font.color.rgb:
                        colors.append(str(run.font.color.rgb))
                except AttributeError:
                    pass
        if idx == 0 and sh.top < 3500000 and t.startswith("内容获客平台"):
            if "FFFFFF" not in colors:
                issues.append(f"S1 主标题非白色 shape{j}")
        if idx == 0 and sh.top > 5500000 and ("+ 全域" in t) and max(sizes or [0]) >= 20:
            issues.append(f"S1 底部误用大标题 shape{j}")
        if FOOTER in t or is_footer_text(t):
            footers.append(j)
        if sh.top > 5500000 and "+ 全域四平台" in t and not is_footer_text(t):
            titles_bottom.append(j)
    if len(footers) > 1:
        issues.append(f"S{idx+1} 重复页脚 shapes={footers}")
    if titles_bottom:
        issues.append(f"S{idx+1} 底部异常标题 shapes={titles_bottom}")
    return issues


def main() -> None:
    src = BASE if BASE.exists() else OUT
    if not src.exists():
        raise SystemExit("缺少 PPT 底稿")
    shutil.copy2(src, OUT)
    prs = Presentation(str(OUT))

    for i, slide in enumerate(prs.slides):
        if i == 0:
            fix_cover(slide)
            fix_footers(slide, i)
        else:
            apply_text_map(slide)
            fix_content_titles(slide)
            fix_footers(slide, i)
        if i == 24:
            fix_slide25(slide)
        if i == 26:
            fix_slide27(slide)

    all_issues = []
    for i, slide in enumerate(prs.slides):
        all_issues.extend(audit(slide, i))

    saved = False
    for path in (OUT, DIR / "内容获客平台_全域四平台实施方案_20260804_无页码.pptx"):
        try:
            prs.save(str(path))
            saved = True
            print(f"已保存：{path}")
            break
        except PermissionError:
            continue
    if not saved:
        alt = DIR / "内容获客平台_全域四平台实施方案_20260804_repaired.pptx"
        prs.save(str(alt))
        print(f"文件占用，另存：{alt}")

    if all_issues:
        print("剩余问题：")
        for x in all_issues:
            print(" ", x)
    else:
        print("逐页自检通过（封面/页脚/重复标题）")


if __name__ == "__main__":
    main()
