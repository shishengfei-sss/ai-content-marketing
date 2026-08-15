#!/usr/bin/env python3
"""更新「内容获客平台_全域四平台实施方案」PPT：融入 AI 内容营销智能体、分期路线、素材合规。"""

from __future__ import annotations

import re
import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

ROOT = Path(__file__).resolve().parents[1]
PPT_PATH = ROOT / "docs/07-对外资料/03-汇报材料/内容获客平台_全域四平台实施方案_20260804.pptx"
BACKUP_PATH = PPT_PATH.with_name(PPT_PATH.stem + "_backup_20260805.pptx")

TOTAL_AFTER = 27
INSERT_AFTER_INDEX = 8  # 0-based：在第 9 页（三层架构）之后插入


def replace_in_slide(slide, mapping: dict[str, str]) -> None:
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = shape.text
        new = text
        for old, val in mapping.items():
            if old in new:
                new = new.replace(old, val)
        if new != text:
            shape.text = new


def replace_all_slides(prs: Presentation, mapping: dict[str, str]) -> None:
    for slide in prs.slides:
        replace_in_slide(slide, mapping)


def renumber_footers(prs: Presentation, total: int) -> None:
    for i, slide in enumerate(prs.slides):
        n = i + 1
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            m = re.fullmatch(r"\d{2}\s*/\s*\d+", shape.text.strip())
            if m:
                shape.text = f"{n:02d} / {total}"


def duplicate_slide(prs: Presentation, source_index: int, insert_index: int):
    """复制幻灯片并插入到指定位置（0-based insert_index）。"""
    source = prs.slides[source_index]
    layout = source.slide_layout
    dest = prs.slides.add_slide(layout)
    for shape in source.shapes:
        newel = deepcopy(shape.element)
        dest.shapes._spTree.insert_element_before(newel, "p:extLst")

    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    new_id = slides[-1]
    xml_slides.remove(new_id)
    xml_slides.insert(insert_index, new_id)
    return prs.slides[insert_index]


def clear_slide_text(slide) -> None:
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.has_text_frame:
            if shape.text.strip() in ("", "›"):
                continue
            # 保留页脚与页眉条：靠位置判断
            if shape.top > 6500000 and len(shape.text) < 80:
                continue
            shape.text = ""


def set_title_block(slide, title: str, subtitle: str) -> None:
    """匹配现有版式：约 top=164592 标题、top=502920 副标题。"""
    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.has_text_frame:
            continue
        if 140000 < shape.top < 220000 and shape.width > 5000000:
            shape.text = title
        elif 450000 < shape.top < 600000 and shape.width > 5000000:
            shape.text = subtitle


def fill_content_slide(slide, blocks: list[tuple[str, str]]) -> None:
    """用左右栏结构填充：左侧标签 + 右侧正文（仿 slide 9 样式）。"""
    rows = [s for s in slide.shapes if hasattr(s, "text") and s.has_text_frame and s.width > 3000000 and s.height < 500000]
    # 找内容区 shape（y > 1M）
    content_shapes = sorted(
        [s for s in slide.shapes if hasattr(s, "text") and s.has_text_frame and s.top > 1000000 and s.left > 900000],
        key=lambda s: (s.top, s.left),
    )
    idx = 0
    for label, body in blocks:
        if idx + 1 < len(content_shapes):
            content_shapes[idx].text = label
            content_shapes[idx + 1].text = body
            idx += 2


def build_new_slide_10(slide) -> None:
    clear_slide_text(slide)
    set_title_block(
        slide,
        "双轮驱动 · AI 内容工厂 + 交易闭环",
        "不只卖课工具 —— 生产内容 → 挂载四平台 → 成交履约 → 持续运营",
    )
    blocks = [
        ("AI 内容工厂（生产层）", "营销创作顾问 · 多平台文案/笔记/视频脚本\n知识库 + 品牌语气 + 合规自检\nPhase 2 起：一键生成商品详情与推广素材"),
        ("内容获客平台（交易层）", "统一商品/订单/支付/权益 · 四平台挂载\n虚拟/服务先闭环 · 课为首发 SKU\nPhase 1 主攻：能挂 · 能卖 · 能履约"),
        ("智营 CRM（运营层，可选）", "活动关联 · 线索/客户跟进 · 成交后经营\n与商城数据互通 · 非强制 Phase 1\n复用已有 v0.6～v1.6 营客能力"),
        ("差异化一句话", "小鹅通/千聊卖「店」—— 我们卖「AI 帮你生产获客内容 + 店里成交」的闭环"),
    ]
    fill_content_slide(slide, blocks)


def build_new_slide_11(slide) -> None:
    clear_slide_text(slide)
    set_title_block(
        slide,
        "内容营销智能体 · 分期能力路线",
        "战略一体、落地分期 —— 不与商城 Phase 1 抢工期",
    )
    blocks = [
        ("Phase 0 ✅ 已有", "多平台文案/笔记/视频脚本 · 活动关联 CRM\nv0.6 顾问 + v1.1 创作增强 · 可演示"),
        ("Phase 1 商城 MVP", "交易闭环优先 · 智能体「复用导出」上架\n不写 AI 配图/成片进硬验收"),
        ("Phase 2 打通", "课纲+详情+卖点一键生成商品 · 创作关联活动/商品\n建议 v1.7 立项"),
        ("Phase 3 图文成品", "AI 商用配图 + 多图 ZIP 可发布包 · 素材授权可追溯\n建议 v1.8"),
        ("Phase 4 短视频", "脚本→TTS→商用授权 BGM→字幕→半自动成片 · 人审后导出\n建议 v1.9"),
        ("Phase 5 效果闭环", "A/B 素材 · 渠道归因 · 发布效果回流 · Multi-Agent\n建议 v2.0+"),
    ]
    fill_content_slide(slide, blocks[:6])


def build_new_slide_12(slide) -> None:
    clear_slide_text(slide)
    set_title_block(
        slide,
        "素材合规 · 图 / 视 / 音三道闸",
        "做得快，更要用得稳 —— 降低公域发布侵权与平台限流风险",
    )
    blocks = [
        ("选源闸", "仅商用授权白名单：图库 / 视频 B-roll / 音乐曲库\nAI 生成走商用套餐 · 用户上传须声明授权"),
        ("成片闸", "每张图/每段视频/每首 BGM 带 license 元数据\n未授权素材不可进入最终导出包"),
        ("发布闸", "文案合规（已有）+ 素材合规（新增）\n人审确认后导出 · 不建议承诺全自动公域发布"),
        ("分阶段开放", "Phase 3：商用配图 · Phase 4：商用 BGM + 半自动成片\n禁止：流行曲随意拼 / 无授权 AI 图默认商用"),
    ]
    fill_content_slide(slide, blocks)


def main() -> None:
    if not PPT_PATH.exists():
        raise SystemExit(f"找不到 PPT：{PPT_PATH}")

    shutil.copy2(PPT_PATH, BACKUP_PATH)
    prs = Presentation(str(PPT_PATH))

    # --- 全文替换 ---
    global_replacements = {
        "24 页": f"{TOTAL_AFTER} 页",
        "/ 24": f"/ {TOTAL_AFTER}",
        "P9–P12": "P9–P15",
        "≈ 7 分钟": "≈ 10 分钟",
        "三层架构 → Phase1 → 四平台 → 店铺归属 A/B": "三层架构 → AI双轮 → 智能体分期 → Phase1 → 四平台",
        "P13–P20": "P16–P23",
        "P21–P24": "P24–P27",
    }
    replace_all_slides(prs, global_replacements)

    # --- 封面 slide 1 ---
    replace_in_slide(
        prs.slides[0],
        {
            "实施方案 · 决策层 24 页": f"实施方案 · 决策层 {TOTAL_AFTER} 页",
            "IP 内容矩阵获客 · 抖音优先": "AI 内容工厂 + IP 矩阵获客",
        },
    )

    # --- slide 2 目录补充 ---
    replace_in_slide(
        prs.slides[1],
        {
            "我们做什么\n---\nP9–P12": f"我们做什么\n---\nP9–P15",
        },
    )

    # --- slide 9 架构页副标题 ---
    replace_in_slide(
        prs.slides[8],
        {
            "公域层 · 内容获客平台核心（多品类） · 经营层可选": "AI 生产层 · 公域层 · 内容获客平台 · 经营层（CRM）",
        },
    )

    # --- 插入 3 页（复制 slide 9 版式）---
    base_index = INSERT_AFTER_INDEX
    for offset, builder in enumerate([build_new_slide_10, build_new_slide_11, build_new_slide_12]):
        new_slide = duplicate_slide(prs, base_index, base_index + 1 + offset)
        builder(new_slide)

    # 插入后原 slide 索引 +3
    def slide_at(old_1based: int):
        return prs.slides[old_1based - 1 + 3]

    # --- slide 20+3=23 差异化 ---
    replace_in_slide(
        slide_at(20),
        {
            "三条护城河": "四条护城河",
            "分阶段开店\n---\n轻店先 · 装修商城后补": "分阶段开店\n---\n轻店先 · 装修商城后补\n---\nAI 内容工厂\n---\n文案→图文→短视频 · 竞品难复制",
        },
    )

    # --- slide 21+3=24 获客 ---
    replace_in_slide(
        slide_at(21),
        {
            "· 内容工厂/AI剪辑不进 Phase 1": "· Phase 2 起：AI 生成详情/推广素材\n· Phase 3/4：商用配图与短视频成片",
            "获客打法 · IP 内容矩阵 → 内容获客平台": "获客打法 · AI 内容工厂 + IP 矩阵 → 内容获客平台",
            "系列知识视频引流 · 店内成交虚拟/服务 · 履约在自有端": "AI 辅助生产内容 → 系列视频引流 → 店内成交 → 自有端履约",
        },
    )

    # --- slide 22+3=25 风险 ---
    replace_in_slide(
        slide_at(22),
        {
            "与千聊同质化": "与千聊同质化 / 素材侵权",
            "多品类虚拟/服务 + 品牌可升级 + IP矩阵": "多品类 + 品牌可升级 + AI内容工厂 + 商用素材合规",
            "Plan B（M6 未达标）": "素材版权风险\n---\n商用白名单曲库/图库 · 授权可追溯 · 人审发布\n禁止无授权流行曲/爬图\n---\nPlan B（M6 未达标）",
        },
    )

    # --- slide 24+3=27 总结 ---
    replace_in_slide(
        slide_at(24),
        {
            "虚拟服务先\n开店能力后补": "AI内容工厂+虚拟服务\n交易闭环优先",
            "IP内容矩阵\n→内容获客平台": "AI生产→IP引流\n→内容获客平台",
            "③ M3 支付硬验收 → Mx 公域首单 → 再补商城开店/第二平台/实物": "③ M3 支付硬验收 → Mx 公域首单\n④ Phase2 智能体×商城打通 → Phase3/4 图文视频",
        },
    )

    renumber_footers(prs, TOTAL_AFTER)
    prs.save(str(PPT_PATH))
    print(f"已更新：{PPT_PATH}")
    print(f"备份：{BACKUP_PATH}")
    print(f"总页数：{len(prs.slides)}")


if __name__ == "__main__":
    main()
