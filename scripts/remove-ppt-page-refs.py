#!/usr/bin/env python3
"""移除 PPT 内章节页码指引（P3–P8）及正文跨页引用（见 P13 等）。"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
PPT = ROOT / "docs/07-对外资料/03-汇报材料/内容获客平台_全域四平台实施方案_20260804_无页码.pptx"

# 目录页独立页码块：P3–P8、P9–P15 …
TOC_PAGE = re.compile(r"^P\d+[–\-]P?\d+$")


def clean_body(text: str) -> str:
    """去掉幻灯片页码引用，保留语义。"""
    t = text
    # 固定短语优先
    replacements = [
        ("公域技术链路 ①抖店 / ②泛知识见 P13，Mx", "公域技术链路 ①抖店 / ②泛知识，Mx"),
        ("（见 P13）", ""),
        ("见 P13，可替代", "可替代"),
        ("（见 P13）。", "。"),
        ("（二选一先通，见 P13）", "（二选一先通）"),
        ("· 规则 P14 · 邀约 P16 · 闭环 P15 · 链路 P13", "· 政策规则 · 邀约对照 · 角色闭环 · 公域链路"),
        ("官方可点击规则链接见 P14。角色闭环见 P15。", "详见政策规则与角色闭环章节。"),
        ("角色分工：个人 IP → 路径 A；平台方 → 链路①抖店证 或 链路②泛知识证（见 P13）。",
         "角色分工：个人 IP → 路径 A；平台方 → 链路①抖店证 或 链路②泛知识证。"),
        ("平台方按 P13 选定链路", "平台方选定链路"),
        ("选定链路 ①或②（P13）+ P15 闭环 + P17 准备清单",
         "选定链路 ①或② + 角色闭环 + 准备清单"),
        ("链路①教培邀约 或 链路②泛知识（见 P13）", "链路①教培邀约 或 链路②泛知识"),
    ]
    for old, new in replacements:
        t = t.replace(old, new)
    # 兜底：去掉残留「见 Pxx」
    t = re.sub(r"\s*见\s*P\d{1,2}\s*", " ", t)
    t = re.sub(r"[（(]\s*P\d{1,2}\s*[）)]", "", t)
    t = re.sub(r"\s+P\d{1,2}\s+闭环", " 闭环", t)
    t = re.sub(r"\s+P\d{1,2}\s+准备", " 准备", t)
    t = re.sub(r" +", " ", t)
    t = re.sub(r" \n", "\n", t)
    t = re.sub(r"\n{3,}", "\n\n", t)
    return t.strip()


def remove_toc_page_blocks(slide) -> int:
    n = 0
    for sh in list(slide.shapes):
        if hasattr(sh, "text") and TOC_PAGE.fullmatch(sh.text.strip()):
            el = sh.element
            el.getparent().remove(el)
            n += 1
    return n


def main() -> None:
    prs = Presentation(str(PPT))
    toc_removed = remove_toc_page_blocks(prs.slides[1])
    body_fixed = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if not hasattr(sh, "text") or not sh.text.strip():
                continue
            if TOC_PAGE.fullmatch(sh.text.strip()):
                continue
            if "P" not in sh.text:
                continue
            # 跳过产品路线图里的 Phase 简写（→P2 / →P3 在 Phase1 后移列表）
            if re.search(r"→\s*P[23]\b", sh.text) and "见 P" not in sh.text and "P1" not in sh.text[:3]:
                if not re.search(r"见\s*P|P\d+–|P\d+·", sh.text):
                    continue
            new = clean_body(sh.text)
            if new != sh.text.strip():
                sh.text = new
                body_fixed += 1
    prs.save(str(PPT))
    print(f"目录页码块移除 {toc_removed} 处；正文跨页引用清理 {body_fixed} 处")
    print(f"已保存：{PPT}")


if __name__ == "__main__":
    main()
